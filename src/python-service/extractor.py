import csv as _csv
import email as _email
import logging
import mimetypes
import os
import subprocess
from pathlib import Path

logger = logging.getLogger(__name__)

STORAGE_ROOT = os.getenv("DMS_STORAGE_PATH", "/app/uploads")
MAX_CONTENT_CHARS = 500_000

_MIME_MAP = {
    ".pdf":  "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".doc":  "application/msword",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".xls":  "application/vnd.ms-excel",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".ppt":  "application/vnd.ms-powerpoint",
    ".txt":  "text/plain",
    ".rtf":  "application/rtf",
    ".odt":  "application/vnd.oasis.opendocument.text",
    ".ods":  "application/vnd.oasis.opendocument.spreadsheet",
    ".odp":  "application/vnd.oasis.opendocument.presentation",
    ".csv":  "text/csv",
    ".html": "text/html",
    ".htm":  "text/html",
    ".xml":  "application/xml",
    ".md":   "text/markdown",
    ".eml":  "message/rfc822",
    ".msg":  "application/vnd.ms-outlook",
}


def resolve_path(file_url: str) -> str:
    return os.path.join(STORAGE_ROOT, file_url.lstrip("/").replace("\\", "/"))


def get_file_info(file_url: str) -> dict:
    if not file_url:
        return {}
    full_path = resolve_path(file_url)
    p = Path(full_path)
    ext = p.suffix.lower()
    size = 0
    try:
        size = p.stat().st_size
    except OSError:
        pass
    return {
        "file_name": p.name,
        "extension": ext,
        "mime_type": _MIME_MAP.get(ext) or mimetypes.guess_type(str(p))[0] or "application/octet-stream",
        "size":      size,
        "path":      file_url,
    }


def extract_text(file_url: str) -> tuple[str, str | None]:
    """Returns (text, error_or_None). Never raises."""
    if not file_url:
        return "", "no file_url provided"
    full_path = resolve_path(file_url)
    if not os.path.isfile(full_path):
        return "", f"file not found: {full_path}"
    ext = Path(full_path).suffix.lower()
    try:
        fn = _EXTRACTORS.get(ext)
        if fn is None:
            return "", f"unsupported extension: {ext}"
        return fn(full_path), None
    except Exception as exc:
        logger.warning("Text extraction failed [%s]: %s", file_url, exc)
        return "", str(exc)


# ── Extractors ────────────────────────────────────────────────

def _pdf(path: str) -> str:
    import fitz
    with fitz.open(path) as doc:
        return "\n".join(page.get_text() for page in doc)


def _docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _doc(path: str) -> str:
    result = subprocess.run(
        ["antiword", path], capture_output=True, timeout=30)
    return result.stdout.decode("utf-8", errors="ignore")


def _xlsx(path: str) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            line = " ".join(str(c) for c in row if c is not None)
            if line.strip():
                parts.append(line)
    return "\n".join(parts)


def _xls(path: str) -> str:
    import xlrd
    wb = xlrd.open_workbook(path)
    parts = []
    for sheet in wb.sheets():
        for row in range(sheet.nrows):
            line = " ".join(str(sheet.cell_value(row, col)) for col in range(sheet.ncols))
            if line.strip():
                parts.append(line)
    return "\n".join(parts)


def _pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
    return "\n".join(parts)


def _ppt(path: str) -> str:
    result = subprocess.run(
        ["catppt", path], capture_output=True, timeout=30)
    return result.stdout.decode("utf-8", errors="ignore")


def _txt(path: str) -> str:
    import chardet
    with open(path, "rb") as f:
        raw = f.read()
    enc = chardet.detect(raw).get("encoding") or "utf-8"
    return raw.decode(enc, errors="ignore")


def _rtf(path: str) -> str:
    from striprtf.striprtf import rtf_to_text
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return rtf_to_text(f.read())


def _odf(path: str) -> str:
    from odf.opendocument import load
    from odf import teletype
    from odf.text import P
    doc = load(path)
    parts = []
    for para in doc.body.getElementsByType(P):
        t = teletype.extractText(para)
        if t.strip():
            parts.append(t)
    return "\n".join(parts)


def _csv(path: str) -> str:
    parts = []
    with open(path, "r", encoding="utf-8", errors="ignore", newline="") as f:
        for row in _csv.reader(f):
            line = " ".join(row)
            if line.strip():
                parts.append(line)
    return "\n".join(parts)


def _html(path: str) -> str:
    from bs4 import BeautifulSoup
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        soup = BeautifulSoup(f.read(), "lxml")
    return soup.get_text(separator="\n", strip=True)


def _xml(path: str) -> str:
    from lxml import etree
    tree = etree.parse(path)
    return " ".join(tree.getroot().itertext()).strip()


def _md(path: str) -> str:
    return _txt(path)


def _eml(path: str) -> str:
    with open(path, "rb") as f:
        msg = _email.message_from_binary_file(f, policy=_email.policy.default)
    parts = [msg.get("subject", ""), msg.get("from", "")]
    if msg.is_multipart():
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                try:
                    parts.append(part.get_content())
                except Exception:
                    pass
    else:
        try:
            parts.append(msg.get_content())
        except Exception:
            pass
    return "\n".join(p for p in parts if p)


def _msg(path: str) -> str:
    import extract_msg
    with extract_msg.Message(path) as msg:
        parts = [msg.subject or "", msg.sender or "", msg.body or ""]
    return "\n".join(p for p in parts if p)


# ── Dispatch table ────────────────────────────────────────────
_EXTRACTORS = {
    ".pdf":  _pdf,
    ".docx": _docx,
    ".doc":  _doc,
    ".xlsx": _xlsx,
    ".xls":  _xls,
    ".pptx": _pptx,
    ".ppt":  _ppt,
    ".txt":  _txt,
    ".rtf":  _rtf,
    ".odt":  _odf,
    ".ods":  _odf,
    ".odp":  _odf,
    ".csv":  _csv,
    ".html": _html,
    ".htm":  _html,
    ".xml":  _xml,
    ".md":   _md,
    ".eml":  _eml,
    ".msg":  _msg,
}
